from __future__ import annotations

import io
from typing import Dict, List

from googleapiclient.discovery import Resource
from googleapiclient.http import MediaIoBaseDownload
from pptx import Presentation

MIME_FOLDER = "application/vnd.google-apps.folder"
MIME_GOOGLE_DOC = "application/vnd.google-apps.document"
MIME_GOOGLE_SLIDES = "application/vnd.google-apps.presentation"
MIME_PDF = "application/pdf"
MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
MIME_PPT = "application/vnd.ms-powerpoint"
MIME_SHORTCUT = "application/vnd.google-apps.shortcut"


def list_study_files(service: Resource, folder_id: str) -> List[Dict]:
    results: List[Dict] = []

    def _recurse(fid: str):
        page_token = None
        while True:
            resp = (
                service.files()
                .list(
                    q=f"'{fid}' in parents and trashed=false",
                    fields=(
                        "nextPageToken, files(id, name, mimeType, "
                        "shortcutDetails(targetId, targetMimeType))"
                    ),
                    pageToken=page_token,
                    pageSize=1000,
                    includeItemsFromAllDrives=True,
                    supportsAllDrives=True,
                )
                .execute()
            )
            for item in resp.get("files", []):
                mt = item.get("mimeType")
                if mt == MIME_FOLDER:
                    _recurse(item["id"])  # type: ignore[index]
                elif mt in (MIME_PDF, MIME_GOOGLE_DOC, MIME_GOOGLE_SLIDES, MIME_PPTX, MIME_PPT):
                    results.append({
                        "id": item["id"],
                        "name": item.get("name", "Untitled"),
                        "mimeType": mt,
                    })
                elif mt == MIME_SHORTCUT:
                    sc = item.get("shortcutDetails", {}) or {}
                    target_id = sc.get("targetId")
                    target_mt = sc.get("targetMimeType")
                    if target_id and target_mt in (MIME_PDF, MIME_GOOGLE_DOC, MIME_GOOGLE_SLIDES, MIME_PPTX, MIME_PPT):
                        results.append({
                            "id": target_id,
                            "name": item.get("name", "Untitled"),
                            "mimeType": target_mt,
                        })
            page_token = resp.get("nextPageToken")
            if not page_token:
                break

    _recurse(folder_id)
    return results


def download_pdf(service: Resource, file_id: str) -> bytes:
    request = service.files().get_media(fileId=file_id)
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return fh.getvalue()


def export_google_doc_as_text(service: Resource, file_id: str) -> str:
    data = (
        service.files()
        .export(fileId=file_id, mimeType="text/plain")
        .execute()
    )
    if isinstance(data, bytes):
        return data.decode("utf-8", errors="ignore")
    # Some clients may return str
    return str(data)


def download_pptx(service: Resource, file_id: str) -> bytes:
    request = service.files().get_media(fileId=file_id)
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return fh.getvalue()


def extract_pptx_text(service: Resource, file_id: str) -> str:
    data = download_pptx(service, file_id)
    prs = Presentation(io.BytesIO(data))
    lines: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            # Regular text frames
            if getattr(shape, "has_text_frame", False):
                tf = shape.text_frame
                if tf is not None:
                    txt = tf.text or ""
                    if txt.strip():
                        for part in txt.splitlines():
                            part = part.strip()
                            if part:
                                lines.append(part)
            # Tables: collect text from cells
            if getattr(shape, "has_table", False):
                tbl = shape.table
                for row in tbl.rows:
                    for cell in row.cells:
                        txt = cell.text or ""
                        if txt.strip():
                            for part in txt.splitlines():
                                part = part.strip()
                                if part:
                                    lines.append(part)
        # Speaker notes, if any
        notes = None
        try:
            if hasattr(slide, "notes_slide") and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
        except Exception:
            notes = None
        if notes:
            for part in notes.splitlines():
                part = part.strip()
                if part:
                    lines.append(part)
    return "\n".join(lines)


def export_google_slides_as_text(service: Resource, file_id: str) -> str:
    data = (
        service.files()
        .export(fileId=file_id, mimeType="text/plain")
        .execute()
    )
    if isinstance(data, bytes):
        return data.decode("utf-8", errors="ignore")
    return str(data)


def export_google_slides_as_pdf(service: Resource, file_id: str) -> bytes:
    data = (
        service.files()
        .export(fileId=file_id, mimeType="application/pdf")
        .execute()
    )
    if isinstance(data, bytes):
        return data
    return str(data).encode("utf-8", errors="ignore")


def list_subfolders(service: Resource, parent_folder_id: str) -> List[Dict]:
    """
    Return a list of subfolders directly under parent_folder_id.
    Each item: { "id": str, "name": str }
    """
    items: List[Dict] = []
    page_token = None
    while True:
        resp = (
            service.files()
            .list(
                q=(
                    f"'{parent_folder_id}' in parents and "
                    f"mimeType='{MIME_FOLDER}' and trashed=false"
                ),
                fields="nextPageToken, files(id, name)",
                pageToken=page_token,
                pageSize=1000,
                includeItemsFromAllDrives=True,
                supportsAllDrives=True,
            )
            .execute()
        )
        for item in resp.get("files", []):
            items.append({"id": item["id"], "name": item.get("name", "Untitled")})
        page_token = resp.get("nextPageToken")
        if not page_token:
            break
    # Sort by name for stable UI
    items.sort(key=lambda x: x.get("name", "").lower())
    return items


def get_folder_name(service: Resource, folder_id: str) -> str:
    meta = service.files().get(fileId=folder_id, fields="id, name", supportsAllDrives=True).execute()
    return meta.get("name", folder_id)


def collect_files_recursively(service: Resource, folder_id: str) -> List[Dict]:
    """
    Recursively walk starting from folder_id and collect all PDF + Google Doc files.
    Uses the existing list_study_files recursion.
    """
    return list_study_files(service, folder_id)
